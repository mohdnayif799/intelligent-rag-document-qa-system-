import io
import os
import concurrent.futures

import numpy as np
import streamlit as st
from PIL import Image

import fitz
import easyocr

from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
)
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

# ── Constants ────────────────────────────────────────────────────────────────
MIN_TEXT_CHARS_PER_PAGE = 20
OCR_RENDER_DPI = 300
OCR_MAX_WORKERS = 4

# Supported extension -> loader strategy
SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".docx", ".doc", ".md", ".pptx"}


# ── OCR helpers ──────────────────────────────────────────────────────────────

@st.cache_resource(show_spinner="Loading OCR engine (first run only)…")
def get_ocr_reader():
    """Load EasyOCR once and cache it. gpu=False works on any machine."""
    return easyocr.Reader(["en"], gpu=False)


def _is_text_sufficient(text: str) -> bool:
    return len(text.strip()) >= MIN_TEXT_CHARS_PER_PAGE


def _ocr_page(args) -> tuple:
    """Render one PDF page to an image and run OCR — thread-safe."""
    pdf_path, page_index = args
    reader = get_ocr_reader()
    doc = fitz.open(pdf_path)
    try:
        page = doc[page_index]
        pix = page.get_pixmap(dpi=OCR_RENDER_DPI)
        image = Image.open(io.BytesIO(pix.tobytes("png")))
        result = reader.readtext(np.array(image), detail=0)
        return page_index, "\n".join(result)
    finally:
        doc.close()


# ── Per-format loaders ───────────────────────────────────────────────────────

def _load_pdf(path: str, progress_callback=None) -> list:
    """
    Load a PDF with per-page OCR fallback.
    Pages with little extractable text (scanned/photographed) are automatically
    routed through EasyOCR so mixed PDFs are handled correctly.
    """
    pages = PyPDFLoader(path).load()
    total = len(pages)

    needs_ocr = [(i, p) for i, p in enumerate(pages) if not _is_text_sufficient(p.page_content)]
    good_pages = {i: p for i, p in enumerate(pages) if _is_text_sufficient(p.page_content)}

    ocr_results = {}
    if needs_ocr:
        ocr_args = [(path, i) for i, _ in needs_ocr]
        with concurrent.futures.ThreadPoolExecutor(max_workers=min(OCR_MAX_WORKERS, len(needs_ocr))) as pool:
            for page_idx, ocr_text in pool.map(_ocr_page, ocr_args):
                ocr_results[page_idx] = ocr_text

    docs = []
    for i, page_doc in enumerate(pages):
        if i in ocr_results:
            docs.append(Document(page_content=ocr_results[i], metadata=page_doc.metadata))
        else:
            docs.append(page_doc)
        if progress_callback:
            progress_callback(i + 1, total)

    if needs_ocr:
        print(f"[INFO] {os.path.basename(path)}: {len(needs_ocr)} page(s) OCR-ed, "
              f"{len(good_pages)} had embedded text.")
    else:
        print(f"[INFO] Loaded {total} page(s) from: {os.path.basename(path)}")
    return docs


def _load_docx(path: str) -> list:
    """Load a DOCX (or DOC) file using Docx2txtLoader."""
    loaded = Docx2txtLoader(path).load()
    # Attach page=0 so metadata is consistent with other loaders
    for doc in loaded:
        doc.metadata.setdefault("page", 0)
    print(f"[INFO] Loaded DOCX: {os.path.basename(path)}")
    return loaded


def _load_text(path: str) -> list:
    """Load a plain-text or Markdown file."""
    loaded = TextLoader(path, encoding="utf-8").load()
    for doc in loaded:
        doc.metadata.setdefault("page", 0)
    print(f"[INFO] Loaded text file: {os.path.basename(path)}")
    return loaded


def _load_pptx(path: str) -> list:
    """
    Extract text from each PowerPoint slide as a separate Document.
    Uses python-pptx directly (lightweight, no heavy unstructured dependency).
    """
    from pptx import Presentation as PptxPresentation  # lazy import

    prs = PptxPresentation(path)
    docs = []
    for slide_num, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        text = "\n".join(parts)
        if text:
            docs.append(Document(
                page_content=text,
                metadata={"source": path, "page": slide_num},
            ))
    print(f"[INFO] Loaded {len(docs)} slide(s) from: {os.path.basename(path)}")
    return docs


# ── Public API ───────────────────────────────────────────────────────────────

def load_documents(file_paths: list, progress_callback=None) -> list:
    """
    Load PDF, DOCX, TXT, MD, and PPTX files from a list of absolute paths.

    - PDFs: fast text extraction with per-page OCR fallback for scanned pages.
    - DOCX/DOC: Docx2txtLoader.
    - TXT/MD: TextLoader.
    - PPTX: python-pptx, one Document per slide.

    progress_callback: optional callable(current, total) for progress bars.
    """
    documents = []

    for path in file_paths:
        ext = os.path.splitext(path)[1].lower()

        if ext == ".pdf":
            documents.extend(_load_pdf(path, progress_callback))

        elif ext in {".docx", ".doc"}:
            documents.extend(_load_docx(path))
            if progress_callback:
                progress_callback(1, 1)

        elif ext in {".txt", ".md"}:
            documents.extend(_load_text(path))
            if progress_callback:
                progress_callback(1, 1)

        elif ext == ".pptx":
            documents.extend(_load_pptx(path))
            if progress_callback:
                progress_callback(1, 1)

        else:
            print(f"[WARNING] Unsupported file type, skipping: {path}")

    return documents


def chunk_documents(documents: list, chunk_size: int = 1000, chunk_overlap: int = 200) -> list:
    """
    Split documents into overlapping chunks.
    Metadata (source, page) is propagated to every chunk automatically
    by RecursiveCharacterTextSplitter.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
        add_start_index=True,   # adds char offset in metadata for better traceability
    )
    chunks = splitter.split_documents(documents)
    print(f"[INFO] Created {len(chunks)} chunks from {len(documents)} document page(s).")
    return chunks
